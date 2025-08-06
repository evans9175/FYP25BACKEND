#!/usr/bin/env python3
"""
Create a sample PowerPoint presentation using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

def create_sample_presentation():
    # Create a presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    title_content_layout = prs.slide_layouts[1]
    section_header_layout = prs.slide_layouts[2]
    two_content_layout = prs.slide_layouts[3]
    comparison_layout = prs.slide_layouts[4]
    title_only_layout = prs.slide_layouts[5]
    blank_layout = prs.slide_layouts[6]
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Sample Presentation"
    subtitle.text = f"Created with Python\n{datetime.datetime.now().strftime('%B %d, %Y')}"
    
    # Slide 2: Introduction
    slide2 = prs.slides.add_slide(title_content_layout)
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Introduction"
    content2.text = "Welcome to this sample presentation!\n\n" \
                   "This presentation demonstrates:\n" \
                   "• Various slide layouts\n" \
                   "• Text formatting\n" \
                   "• Shapes and graphics\n" \
                   "• Charts and tables"
    
    # Slide 3: Section Header
    slide3 = prs.slides.add_slide(section_header_layout)
    title3 = slide3.shapes.title
    subtitle3 = slide3.placeholders[1]
    
    title3.text = "Key Features"
    subtitle3.text = "Exploring PowerPoint capabilities with Python"
    
    # Slide 4: Two Content Layout
    slide4 = prs.slides.add_slide(two_content_layout)
    title4 = slide4.shapes.title
    left_content = slide4.placeholders[1]
    right_content = slide4.placeholders[2]
    
    title4.text = "Benefits of Automated Presentations"
    
    left_content.text = "Efficiency\n\n" \
                       "• Save time\n" \
                       "• Reduce errors\n" \
                       "• Consistent formatting\n" \
                       "• Easy updates"
    
    right_content.text = "Flexibility\n\n" \
                        "• Data-driven content\n" \
                        "• Dynamic generation\n" \
                        "• Version control\n" \
                        "• Batch processing"
    
    # Slide 5: Content with Shapes
    slide5 = prs.slides.add_slide(title_only_layout)
    title5 = slide5.shapes.title
    title5.text = "Visual Elements"
    
    # Add shapes
    left = Inches(1)
    top = Inches(2)
    width = Inches(2)
    height = Inches(2)
    
    # Rectangle
    shape1 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
    shape1.text = "Rectangle"
    
    # Circle
    shape2 = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(2.5), top, width, height
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(112, 173, 71)  # Green
    shape2.text = "Circle"
    
    # Arrow
    shape3 = slide5.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left + Inches(5), top, width, height
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 192, 0)  # Orange
    shape3.text = "Arrow"
    
    # Slide 6: Table
    slide6 = prs.slides.add_slide(title_only_layout)
    title6 = slide6.shapes.title
    title6.text = "Data Table Example"
    
    # Define table data
    rows = 4
    cols = 3
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(6)
    height = Inches(3)
    
    table = slide6.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    
    # Add header row
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Q1 Results"
    table.cell(0, 2).text = "Q2 Results"
    
    # Add data
    data = [
        ["Product A", "85%", "92%"],
        ["Product B", "78%", "81%"],
        ["Product C", "91%", "88%"]
    ]
    
    for i, row_data in enumerate(data, 1):
        for j, cell_data in enumerate(row_data):
            table.cell(i, j).text = cell_data
    
    # Style the header row
    for i in range(cols):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
    
    # Slide 7: Conclusion
    slide7 = prs.slides.add_slide(title_content_layout)
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Conclusion"
    content7.text = "Key Takeaways:\n\n" \
                   "✓ Python-pptx enables automated presentation creation\n" \
                   "✓ Supports various layouts and formatting options\n" \
                   "✓ Can include shapes, tables, and other visual elements\n" \
                   "✓ Perfect for data-driven presentations\n\n" \
                   "Thank you!"
    
    # Slide 8: End slide
    slide8 = prs.slides.add_slide(title_slide_layout)
    title8 = slide8.shapes.title
    subtitle8 = slide8.placeholders[1]
    
    title8.text = "Questions?"
    subtitle8.text = "Thank you for your attention"
    
    # Save the presentation
    prs.save('sample_presentation.pptx')
    print("Presentation created successfully!")
    print("File saved as: sample_presentation.pptx")

if __name__ == "__main__":
    create_sample_presentation()